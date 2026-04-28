"""File content extraction — extract text from various formats as Markdown."""

from __future__ import annotations

import base64
import logging
from pathlib import Path
from typing import Callable

logger = logging.getLogger('nowork')


def extract_text(file_path: str | Path) -> str | None:
    """Extract text content from a file, returned as Markdown.

    Returns:
        Extracted text, or None for unsupported file types.
    """
    path = Path(file_path)
    if not path.exists():
        return None

    ext = path.suffix.lower()

    extractors: dict[str, Callable[[Path], str | None]] = {
        # Plain text
        '.md': _read_text, '.txt': _read_text,
        '.py': _read_text, '.js': _read_text, '.ts': _read_text,
        '.json': _read_text, '.yaml': _read_text, '.yml': _read_text,
        '.toml': _read_text, '.csv': _read_text, '.xml': _read_text,
        '.html': _read_text, '.css': _read_text, '.sql': _read_text,
        '.sh': _read_text, '.bat': _read_text,
        '.go': _read_text, '.rs': _read_text, '.java': _read_text,
        '.c': _read_text, '.cpp': _read_text, '.h': _read_text, '.hpp': _read_text,
        '.rb': _read_text, '.php': _read_text, '.swift': _read_text,
        '.kt': _read_text, '.r': _read_text, '.R': _read_text,

        # Documents
        '.pdf': _extract_pdf,
        '.docx': _extract_docx,
        '.pptx': _extract_pptx,
        '.xlsx': _extract_xlsx,

        # Images
        '.png': _extract_image,
        '.jpg': _extract_image, '.jpeg': _extract_image,
        '.gif': _extract_image, '.webp': _extract_image,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        logger.debug('Unsupported file type: %s', ext)
        return None

    try:
        return extractor(path)
    except Exception as e:
        logger.warning('Failed to extract text from %s: %s', path, e)
        return None


def _read_text(path: Path) -> str:
    return path.read_text(encoding='utf-8', errors='replace')


def _extract_pdf(path: Path) -> str | None:
    """Extract PDF text via PyMuPDF."""
    try:
        import fitz
    except ImportError:
        logger.warning('PyMuPDF (fitz) not installed, cannot extract PDF')
        return None

    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    doc.close()
    return '\n\n'.join(pages) if pages else None


def _extract_docx(path: Path) -> str | None:
    """Extract DOCX text preserving heading hierarchy via python-docx."""
    try:
        from docx import Document
    except ImportError:
        logger.warning('python-docx not installed, cannot extract DOCX')
        return None

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        style = para.style.name or ''
        text = para.text.strip()
        if not text:
            continue
        if 'Heading' in style:
            level = int(''.join(filter(str.isdigit, style)) or '1')
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)
    return '\n\n'.join(parts) if parts else None


def _extract_xlsx(path: Path) -> str | None:
    """Extract XLSX and convert to Markdown tables via openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning('openpyxl not installed, cannot extract XLSX')
        return None

    wb = load_workbook(str(path), read_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## {ws.title}")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else '' for c in row]
            line = '| ' + ' | '.join(cells) + ' |'
            parts.append(line)
            if i == 0:
                parts.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
    wb.close()
    return '\n'.join(parts) if parts else None


def _extract_pptx(path: Path) -> str | None:
    """Extract PPTX text via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning('python-pptx not installed, cannot extract PPTX')
        return None

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
    return '\n\n'.join(parts) if parts else None


def _extract_image(path: Path) -> str:
    """Return base64-encoded image data for multimodal agents."""
    data = path.read_bytes()
    b64 = base64.b64encode(data).decode()
    ext = path.suffix.lower().lstrip('.')
    return f"[Image: {path.name}]\n(data:image/{ext};base64,{b64})"
